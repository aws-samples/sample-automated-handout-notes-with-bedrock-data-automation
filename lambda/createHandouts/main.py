import boto3
from pptx import Presentation
from pptx.util import Inches
import os
import json
import logging
from botocore.exceptions import ClientError

# Set up logging
logger = logging.getLogger()
logger.setLevel(logging.INFO)

s3 = boto3.client('s3')

def validate_input(event):
    """Validate the input event structure"""
    if not event.get("Payload"):
        raise ValueError("Missing Payload in event")
    
    fixed_event = event["Payload"]
    if len(fixed_event) < 2:
        raise ValueError("Invalid Payload structure")
    
    if not fixed_event[1].get("outputbucket") or not fixed_event[1].get("outputkey"):
        raise ValueError("Missing outputbucket or outputkey in event")
    
    return fixed_event

def lambda_handler(event, context):
    try:
        # Validate input
        fixed_event = validate_input(event)
        
        # Define bucket and file paths    
        bucket_name = fixed_event[1]["outputbucket"]
        images_folder = fixed_event[1]["outputkey"]
        output_key = f"handouts/{images_folder.strip('/').split('screenshots_')[1]}_lecture_notes.pptx"

        # Temporary storage for Lambda
        tmp_dir = '/tmp/'
        os.makedirs(tmp_dir, exist_ok=True)
        
        # List all objects in the images folder
        try:
            response = s3.list_objects_v2(Bucket=bucket_name, Prefix=images_folder)
            images = response.get('Contents', [])
            if not images:
                raise ValueError(f"No images found in bucket {bucket_name} with prefix {images_folder}")
        except ClientError as e:
            logger.error(f"Error accessing S3 bucket: {str(e)}")
            raise
        
        # Filter out just the filenames for images
        image_files = [img['Key'] for img in images if img['Key'] != images_folder]
        
        # Get the list of transcription segments
        transcription_segments = fixed_event[0]
        
        # Validate matching counts
        num_images = len(image_files)
        num_transcripts = len(transcription_segments)
        
        if num_images != num_transcripts:
            raise ValueError(f"Number of images ({num_images}) does not match number of transcripts ({num_transcripts})")
        
        # Sort images
        image_files.sort()
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = int(12192000)
        prs.slide_height = int(6858000)
        slide_width = prs.slide_width
        
        # Process images and transcriptions
        for i in range(num_images):
            try:
                # Download image
                image_key = image_files[i]
                image_path = os.path.join(tmp_dir, os.path.basename(image_key))
                s3.download_file(bucket_name, image_key, image_path)
                
                # Add slide and image
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.add_picture(image_path, 0, 0, width=slide_width)
                
                # Add transcription to notes
                notes_slide = slide.notes_slide
                transcription_text = transcription_segments[i].get('transcript', '')
                if not transcription_text:
                    logger.warning(f"Empty transcription for slide {i+1}")
                notes_slide.notes_text_frame.text = transcription_text
                
                # Clean up temporary file
                os.remove(image_path)
                
            except Exception as e:
                logger.error(f"Error processing slide {i+1}: {str(e)}")
                raise
        
        # Save and upload PowerPoint
        try:
            pptx_path = os.path.join(tmp_dir, "lecture_notes.pptx")
            prs.save(pptx_path)
            
            s3.upload_file(pptx_path, bucket_name, output_key)
            
            # Clean up
            os.remove(pptx_path)
            
        except Exception as e:
            logger.error(f"Error saving or uploading presentation: {str(e)}")
            raise
        
        return {
            "statusCode": 200,
            "message": "PPTX presentation generated successfully!",
            "pptx_url": f"s3://{bucket_name}/{output_key}"
        }
        
    except ValueError as e:
        logger.error(f"Validation error: {str(e)}")
        return {
            "statusCode": 400,
            "message": f"Invalid input: {str(e)}"
        }
    except ClientError as e:
        logger.error(f"AWS error: {str(e)}")
        return {
            "statusCode": 500,
            "message": f"AWS service error: {str(e)}"
        }
    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")
        return {
            "statusCode": 500,
            "message": f"Internal server error: {str(e)}"
        }
